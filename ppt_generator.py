from copy import deepcopy

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TITLE_SIZE = 44
SUBTITLE_SIZE = 28
BODY_SIZE = 18
MARGIN = Inches(1.0)

THEMES = {
    "modern": {
        "background": "#EEF2FF",
        "surface": "#FFFFFF",
        "text": "#0F172A",
        "muted": "#475569",
        "primary": "#4F46E5",
        "secondary": "#6366F1",
        "accent": "#22C55E",
        "line": "#C7D2FE",
    },
    "corporate": {
        "background": "#F8FAFC",
        "surface": "#FFFFFF",
        "text": "#0B1F3A",
        "muted": "#334155",
        "primary": "#1E3A8A",
        "secondary": "#1D4ED8",
        "accent": "#0EA5E9",
        "line": "#BFDBFE",
    },
    "academic": {
        "background": "#F8FAFC",
        "surface": "#FFFFFF",
        "text": "#0F172A",
        "muted": "#334155",
        "primary": "#0F766E",
        "secondary": "#14B8A6",
        "accent": "#F59E0B",
        "line": "#CCFBF1",
    },
    "dark": {
        "background": "#0B1220",
        "surface": "#111827",
        "text": "#E2E8F0",
        "muted": "#94A3B8",
        "primary": "#22D3EE",
        "secondary": "#0EA5E9",
        "accent": "#A78BFA",
        "line": "#1F2937",
    },
}

ICON_MAP = {
    "rocket": "🚀",
    "chart": "📊",
    "book": "📘",
    "ai": "🤖",
    "briefcase": "💼",
    "clipboard": "📋",
    "columns": "🧩",
    "scale": "⚖️",
    "timeline": "🕒",
    "target": "🎯",
    "check": "✅",
    "spark": "✨",
    "handshake": "🤝",
}


def _hex_to_rgb(value: str) -> RGBColor:
    clean = (value or "").strip().lstrip("#")
    if len(clean) == 3:
        clean = "".join(ch * 2 for ch in clean)
    if len(clean) != 6:
        clean = "4F46E5"
    try:
        return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))
    except ValueError:
        return RGBColor(79, 70, 229)


def _palette(theme: str) -> dict:
    safe = (theme or "modern").lower()
    return deepcopy(THEMES.get(safe, THEMES["modern"]))


def _palette_for_slide(base_palette: dict, slide_data: dict) -> dict:
    local = deepcopy(base_palette)
    design = slide_data.get("design") if isinstance(slide_data, dict) else {}
    if isinstance(design, dict):
        accent_color = design.get("accent_color")
        if isinstance(accent_color, str) and accent_color.startswith("#"):
            local["primary"] = accent_color
            local["accent"] = accent_color
    return local


def _add_shape_rect(slide, left, top, width, height, fill_hex, line_hex=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = _hex_to_rgb(line_hex)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, size, color_hex, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    para = frame.paragraphs[0]
    para.text = text
    para.alignment = align
    para.font.name = "Calibri"
    para.font.size = Pt(size)
    para.font.bold = bold
    para.font.color.rgb = _hex_to_rgb(color_hex)
    return box


def _add_bullets(slide, left, top, width, height, bullets, palette, icon="•"):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()

    items = (bullets or [])[:5]
    if not items:
        items = ["Clear objective", "Practical recommendation", "Expected business outcome"]

    for idx, bullet in enumerate(items):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.text = f"{icon} {bullet}"
        para.font.name = "Calibri"
        para.font.size = Pt(BODY_SIZE)
        para.font.bold = False
        para.font.color.rgb = _hex_to_rgb(palette["text"])
        para.level = 0
        para.space_after = Pt(10)


def _add_numbered_points(slide, left, top, width, height, points, palette):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()

    items = (points or [])[:6]
    if not items:
        items = ["Context and constraints", "Strategic options", "Execution plan", "Expected impact"]

    for idx, point in enumerate(items, start=1):
        para = frame.paragraphs[0] if idx == 1 else frame.add_paragraph()
        para.text = f"{idx}. {point}"
        para.font.name = "Calibri"
        para.font.size = Pt(BODY_SIZE)
        para.font.bold = False
        para.font.color.rgb = _hex_to_rgb(palette["text"])
        para.level = 0
        para.space_after = Pt(10)


def _normalize_chart_points(value):
    if isinstance(value, list):
        items = value
    elif isinstance(value, dict):
        items = [{"label": k, "value": v} for k, v in value.items()]
    else:
        items = []

    points = []
    for idx, item in enumerate(items, start=1):
        if isinstance(item, dict):
            label = str(item.get("label") or item.get("name") or "Metric").strip()
            numeric_raw = item.get("value") or item.get("score") or item.get("amount") or 0
        else:
            label = str(item).strip()
            numeric_raw = 0
        try:
            value_num = int(float(str(numeric_raw).replace("%", "").strip()))
        except Exception:
            value_num = 25 + idx * 12
        points.append({"label": label[:20] or "Metric", "value": max(min(value_num, 100), 10)})
        if len(points) >= 5:
            break

    if points and len({point["value"] for point in points}) == 1:
        for idx, point in enumerate(points, start=1):
            point["value"] = max(min(point["value"] + idx * 4, 100), 10)
    return points


def _style_chart(chart, chart_type: str, palette: dict):
    try:
        chart.chart_style = 10
    except Exception:
        pass

    try:
        if chart_type != "pie":
            chart.has_legend = False
        else:
            chart.has_legend = True
            chart.legend.include_in_layout = False
    except Exception:
        pass

    try:
        series = chart.series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _hex_to_rgb(palette["primary"])
        if chart_type == "trend":
            series.format.line.color.rgb = _hex_to_rgb(palette["primary"])
    except Exception:
        pass

    try:
        plot = chart.plots[0]
        if chart_type != "pie":
            plot.has_data_labels = True
            plot.data_labels.number_format = "0"
    except Exception:
        pass

    try:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = _hex_to_rgb(palette["line"])
    except Exception:
        pass


def _slide_background(prs, slide, palette, gradient=False):
    _add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, palette["background"])
    if gradient:
        overlay = _add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height / 2, palette["secondary"])
        overlay.fill.transparency = 0.85


def _chrome(prs, slide, palette):
    _add_shape_rect(slide, Inches(0), Inches(0), Inches(0.22), prs.slide_height, palette["primary"])
    _add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), palette["secondary"])
    _add_shape_rect(slide, Inches(0), prs.slide_height - Inches(0.16), prs.slide_width, Inches(0.16), palette["accent"])
    _add_shape_rect(slide, Inches(0.25), Inches(0.68), prs.slide_width - Inches(0.5), Inches(0.02), palette["line"])


def _icon(icon_name: str) -> str:
    key = (icon_name or "spark").lower()
    return ICON_MAP.get(key, "✨")


def _prepare_slide(prs, slide_data, palette, gradient=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_background(prs, slide, palette, gradient=gradient)
    _chrome(prs, slide, palette)
    return slide


def create_title_slide(prs, slide_data, palette):
    slide = _prepare_slide(prs, slide_data, palette, gradient=True)
    icon_name = ((slide_data.get("design") or {}).get("icon") if isinstance(slide_data.get("design"), dict) else "rocket")
    _add_textbox(slide, Inches(1.1), Inches(1.2), Inches(11.0), Inches(0.9), _icon(icon_name), 38, palette["accent"], bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(
        slide,
        Inches(1.0),
        Inches(2.0),
        Inches(11.3),
        Inches(1.3),
        slide_data.get("title", "Presentation Title"),
        TITLE_SIZE,
        palette["text"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(1.2),
        Inches(3.4),
        Inches(10.8),
        Inches(1.0),
        slide_data.get("subtitle", "Professional presentation"),
        SUBTITLE_SIZE,
        palette["muted"],
        bold=False,
        align=PP_ALIGN.CENTER,
    )

    insight_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(4.75), Inches(9.3), Inches(1.05))
    insight_panel.fill.solid()
    insight_panel.fill.fore_color.rgb = _hex_to_rgb(palette["surface"])
    insight_panel.line.color.rgb = _hex_to_rgb(palette["line"])

    title_points = slide_data.get("bullets") if isinstance(slide_data.get("bullets"), list) else []
    strapline = " | ".join([str(item).strip() for item in title_points[:3] if str(item).strip()])
    if not strapline:
        strapline = "Executive overview • Strategic recommendations • Measurable outcomes"
    _add_textbox(slide, Inches(2.25), Inches(5.05), Inches(8.8), Inches(0.55), strapline, 16, palette["muted"], align=PP_ALIGN.CENTER)


def create_section_slide(prs, slide_data, palette):
    slide = _prepare_slide(prs, slide_data, palette, gradient=True)
    icon_name = ((slide_data.get("design") or {}).get("icon") if isinstance(slide_data.get("design"), dict) else "spark")
    _add_textbox(slide, Inches(1.1), Inches(2.0), Inches(11.0), Inches(1.0), _icon(icon_name), 34, palette["accent"], bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.3), slide_data.get("title", "Section"), 40, palette["text"], bold=True, align=PP_ALIGN.CENTER)
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        _add_textbox(slide, Inches(1.1), Inches(4.2), Inches(11.0), Inches(0.9), subtitle, 22, palette["muted"], align=PP_ALIGN.CENTER)

    section_points = (slide_data.get("bullets") or [])[:3]
    if section_points:
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(5.05), Inches(9.2), Inches(1.35))
        panel.fill.solid()
        panel.fill.fore_color.rgb = _hex_to_rgb(palette["surface"])
        panel.line.color.rgb = _hex_to_rgb(palette["line"])
        _add_bullets(slide, Inches(2.25), Inches(5.2), Inches(8.7), Inches(1.0), section_points, palette, icon="•")


def create_agenda_slide(prs, slide_data, palette):
    slide = _prepare_slide(prs, slide_data, palette, gradient=True)
    _add_textbox(slide, MARGIN, Inches(0.9), Inches(10.8), Inches(1.0), slide_data.get("title", "Agenda"), 38, palette["text"], bold=True)
    subtitle = slide_data.get("subtitle", "Session roadmap")
    _add_textbox(slide, MARGIN, Inches(1.75), Inches(10.8), Inches(0.8), subtitle, 20, palette["muted"])

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.6), Inches(11.0), Inches(3.8))
    panel.fill.solid()
    panel.fill.fore_color.rgb = _hex_to_rgb(palette["surface"])
    panel.line.color.rgb = _hex_to_rgb(palette["line"])

    _add_numbered_points(slide, Inches(1.35), Inches(2.95), Inches(10.3), Inches(3.1), slide_data.get("bullets") or [], palette)


def create_bullet_slide(prs, slide_data, palette):
    slide = _prepare_slide(prs, slide_data, palette)
    icon_name = ((slide_data.get("design") or {}).get("icon") if isinstance(slide_data.get("design"), dict) else "check")

    _add_textbox(slide, MARGIN, Inches(0.9), Inches(8.5), Inches(1.0), slide_data.get("title", "Key Points"), 36, palette["text"], bold=True)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.1), Inches(0.9), Inches(2.0), Inches(0.9))
    badge.fill.solid()
    badge.fill.fore_color.rgb = _hex_to_rgb(palette["primary"])
    badge.line.fill.background()
    badge.text_frame.text = f"{_icon(icon_name)} Focus"
    badge.text_frame.paragraphs[0].font.name = "Calibri"
    badge.text_frame.paragraphs[0].font.size = Pt(16)
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(1.9), Inches(11.0), Inches(4.7))
    panel.fill.solid()
    panel.fill.fore_color.rgb = _hex_to_rgb(palette["surface"])
    panel.line.color.rgb = _hex_to_rgb(palette["line"])

    _add_bullets(slide, Inches(1.25), Inches(2.25), Inches(10.5), Inches(4.0), slide_data.get("bullets") or [], palette, icon="•")


def create_two_column_slide(prs, slide_data, palette):
    slide = _prepare_slide(prs, slide_data, palette)
    _add_textbox(slide, MARGIN, Inches(0.9), Inches(10.5), Inches(1.0), slide_data.get("title", "Two Column"), 34, palette["text"], bold=True)

    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(1.9), Inches(5.3), Inches(4.7))
    left.fill.solid()
    left.fill.fore_color.rgb = _hex_to_rgb(palette["surface"])
    left.line.color.rgb = _hex_to_rgb(palette["line"])

    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(1.9), Inches(5.3), Inches(4.7))
    right.fill.solid()
    right.fill.fore_color.rgb = _hex_to_rgb(palette["surface"])
    right.line.color.rgb = _hex_to_rgb(palette["line"])

    _add_textbox(slide, Inches(1.2), Inches(2.1), Inches(4.6), Inches(0.7), slide_data.get("left_title", "Left"), 24, palette["primary"], bold=True)
    _add_bullets(slide, Inches(1.2), Inches(2.85), Inches(4.6), Inches(3.4), slide_data.get("left_points") or [], palette)

    _add_textbox(slide, Inches(6.3), Inches(2.1), Inches(4.6), Inches(0.7), slide_data.get("right_title", "Right"), 24, palette["primary"], bold=True)
    _add_bullets(slide, Inches(6.3), Inches(2.85), Inches(4.6), Inches(3.4), slide_data.get("right_points") or [], palette)


def create_comparison_slide(prs, slide_data, palette):
    slide = _prepare_slide(prs, slide_data, palette)
    _add_textbox(slide, MARGIN, Inches(0.9), Inches(10.5), Inches(1.0), slide_data.get("title", "Comparison"), 34, palette["text"], bold=True)

    _add_shape_rect(slide, Inches(6.35), Inches(1.9), Inches(0.06), Inches(4.7), palette["accent"])

    _add_textbox(slide, Inches(1.0), Inches(2.0), Inches(5.0), Inches(0.8), slide_data.get("left_title", "Option A"), 24, palette["primary"], bold=True, align=PP_ALIGN.CENTER)
    _add_bullets(slide, Inches(1.0), Inches(2.85), Inches(5.0), Inches(3.5), slide_data.get("left_points") or [], palette, icon="◦")

    _add_textbox(slide, Inches(6.7), Inches(2.0), Inches(5.0), Inches(0.8), slide_data.get("right_title", "Option B"), 24, palette["primary"], bold=True, align=PP_ALIGN.CENTER)
    _add_bullets(slide, Inches(6.7), Inches(2.85), Inches(5.0), Inches(3.5), slide_data.get("right_points") or [], palette, icon="◦")


def create_timeline_slide(prs, slide_data, palette):
    slide = _prepare_slide(prs, slide_data, palette)
    _add_textbox(slide, MARGIN, Inches(0.9), Inches(10.5), Inches(1.0), slide_data.get("title", "Timeline"), 34, palette["text"], bold=True)

    _add_shape_rect(slide, Inches(1.4), Inches(3.4), Inches(10.0), Inches(0.06), palette["primary"])
    milestones = (slide_data.get("milestones") or [])[:5]
    if not milestones:
        milestones = ["Plan", "Build", "Launch"]

    total = len(milestones)
    spacing = 9.6 / max(total - 1, 1)
    for idx, milestone in enumerate(milestones):
        x = 1.4 + idx * spacing
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(3.15), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = _hex_to_rgb(palette["accent"])
        circle.line.fill.background()
        _add_textbox(slide, Inches(x - 0.45), Inches(3.75), Inches(1.4), Inches(1.2), milestone, BODY_SIZE, palette["text"], align=PP_ALIGN.CENTER)


def create_chart_slide(prs, slide_data, palette):
    slide = _prepare_slide(prs, slide_data, palette)
    _add_textbox(slide, MARGIN, Inches(0.9), Inches(10.5), Inches(1.0), slide_data.get("title", "Data Insights"), 34, palette["text"], bold=True)
    chart_title = slide_data.get("chart_title") or slide_data.get("subtitle") or "Performance by pillar"
    _add_textbox(slide, MARGIN, Inches(1.75), Inches(10.5), Inches(0.7), chart_title, 18, palette["muted"])

    chart_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.35), Inches(11.2), Inches(4.0))
    chart_panel.fill.solid()
    chart_panel.fill.fore_color.rgb = _hex_to_rgb(palette["surface"])
    chart_panel.line.color.rgb = _hex_to_rgb(palette["line"])

    points = _normalize_chart_points(slide_data.get("chart_points"))
    if not points:
        points = [
            {"label": "Metric A", "value": 65},
            {"label": "Metric B", "value": 72},
            {"label": "Metric C", "value": 58},
            {"label": "Metric D", "value": 84},
        ]
    chart_data = CategoryChartData()
    chart_data.categories = [point["label"] for point in points]
    chart_data.add_series("Value", [point["value"] for point in points])

    chart_type = (slide_data.get("chart_type") or "bar").lower()
    if chart_type == "pie":
        ppt_chart_type = XL_CHART_TYPE.PIE
    elif chart_type == "trend":
        ppt_chart_type = XL_CHART_TYPE.LINE_MARKERS
    else:
        ppt_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

    chart = slide.shapes.add_chart(
        ppt_chart_type,
        Inches(1.35),
        Inches(2.65),
        Inches(10.5),
        Inches(3.45),
        chart_data,
    ).chart

    _style_chart(chart, chart_type, palette)
    try:
        chart.chart_title.has_text_frame = False
    except Exception:
        pass


def create_summary_slide(prs, slide_data, palette):
    slide = _prepare_slide(prs, slide_data, palette)
    _add_textbox(slide, MARGIN, Inches(0.9), Inches(10.5), Inches(1.0), slide_data.get("title", "Summary"), 34, palette["text"], bold=True)

    highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.9), Inches(11.1), Inches(1.0))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = _hex_to_rgb(palette["primary"])
    highlight.line.fill.background()
    highlight.text_frame.text = slide_data.get("summary_banner", "Key takeaways and next actions")
    highlight.text_frame.paragraphs[0].font.name = "Calibri"
    highlight.text_frame.paragraphs[0].font.bold = True
    highlight.text_frame.paragraphs[0].font.size = Pt(20)
    highlight.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.1), Inches(11.1), Inches(3.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = _hex_to_rgb(palette["surface"])
    panel.line.color.rgb = _hex_to_rgb(palette["line"])

    _add_bullets(slide, Inches(1.3), Inches(3.35), Inches(10.5), Inches(2.8), slide_data.get("bullets") or [], palette, icon="✓")


def create_conclusion_slide(prs, slide_data, palette):
    conclusion_payload = deepcopy(slide_data if isinstance(slide_data, dict) else {})
    conclusion_payload.setdefault("title", "Conclusion")
    conclusion_payload.setdefault("summary_banner", "Final recommendations and next steps")
    create_summary_slide(prs, conclusion_payload, palette)


def create_thank_you_slide(prs, slide_data, palette):
    slide = _prepare_slide(prs, slide_data, palette, gradient=True)
    icon_name = ((slide_data.get("design") or {}).get("icon") if isinstance(slide_data.get("design"), dict) else "handshake")
    _add_textbox(slide, Inches(1.1), Inches(2.0), Inches(11.0), Inches(1.1), _icon(icon_name), 38, palette["accent"], bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.0), Inches(3.0), Inches(11.3), Inches(1.2), slide_data.get("title", "Thank You"), 42, palette["text"], bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.9), slide_data.get("subtitle", "Questions and discussion"), 24, palette["muted"], align=PP_ALIGN.CENTER)

    closing_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(5.0), Inches(9.2), Inches(1.15))
    closing_panel.fill.solid()
    closing_panel.fill.fore_color.rgb = _hex_to_rgb(palette["surface"])
    closing_panel.line.color.rgb = _hex_to_rgb(palette["line"])
    _add_textbox(slide, Inches(2.3), Inches(5.3), Inches(8.6), Inches(0.55), "Let us align on next steps and ownership.", 16, palette["muted"], align=PP_ALIGN.CENTER)


def _render_slide(prs, slide_data, deck_palette):
    palette = _palette_for_slide(deck_palette, slide_data)
    slide_type = (slide_data.get("slide_type") or "bullet").lower()

    if slide_type == "title":
        create_title_slide(prs, slide_data, palette)
    elif slide_type == "agenda":
        create_agenda_slide(prs, slide_data, palette)
    elif slide_type == "section":
        create_section_slide(prs, slide_data, palette)
    elif slide_type == "two-column":
        create_two_column_slide(prs, slide_data, palette)
    elif slide_type == "comparison":
        create_comparison_slide(prs, slide_data, palette)
    elif slide_type == "chart":
        create_chart_slide(prs, slide_data, palette)
    elif slide_type == "timeline":
        create_timeline_slide(prs, slide_data, palette)
    elif slide_type in {"summary", "conclusion"}:
        if slide_type == "conclusion":
            create_conclusion_slide(prs, slide_data, palette)
        else:
            create_summary_slide(prs, slide_data, palette)
    elif slide_type == "thank-you":
        create_thank_you_slide(prs, slide_data, palette)
    else:
        create_bullet_slide(prs, slide_data, palette)


def build_pptx_file(generated: dict, file_path: str):
    prs = Presentation()
    deck_theme = (generated.get("theme") or "modern").lower()
    deck_palette = _palette(deck_theme)

    slides = generated.get("slides") if isinstance(generated.get("slides"), list) else []
    if not slides:
        slides = [
            {
                "slide_type": "title",
                "title": generated.get("title", "Presentation"),
                "subtitle": "Professional deck",
                "layout": "centered",
                "design": {"background": "gradient", "accent_color": deck_palette["primary"], "icon": "rocket", "shape": "rounded"},
            },
            {
                "slide_type": "thank-you",
                "title": "Thank You",
                "subtitle": "Q&A",
                "layout": "centered",
                "design": {"background": "gradient", "accent_color": deck_palette["primary"], "icon": "handshake", "shape": "rounded"},
            },
        ]

    for slide_data in slides:
        _render_slide(prs, slide_data if isinstance(slide_data, dict) else {}, deck_palette)

    prs.save(file_path)
